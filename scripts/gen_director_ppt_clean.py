from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


desktop = Path.home() / 'Desktop'
root = Path(__file__).resolve().parent.parent
output_path = desktop / 'Presentation_SAFE_CONGO_Directeur_Demain.pptx'
image_path = root / 'models' / 'evaluation' / 'confusion_matrices_combined.png'

PRIMARY = RGBColor(10, 95, 171)
SECONDARY = RGBColor(73, 172, 239)
TEXT = RGBColor(31, 41, 55)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(217, 119, 6)
GOOD = RGBColor(5, 150, 105)
BG = RGBColor(245, 248, 252)


def background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.5))
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()


def title(slide, text, subtext=''):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.75), Inches(12), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    if subtext:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.35), Inches(12), Inches(0.35))
        p2 = sub.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtext
        r2.font.size = Pt(11)
        r2.font.color.rgb = TEXT


def bullets(slide, items, left, top, width, height, size=17):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.bullet = True
        p.space_after = Pt(7)


def card(slide, heading, value, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = heading
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(20)
    r2.font.bold = True
    r2.font.color.rgb = WHITE


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, 'SAFE CONGO - pipeline data et performances du modele', 'Presentation pour direction | 6 slides maximum')
bullets(slide, [
    'Point de depart : un premier dataset 2023 seul, insuffisant pour stabiliser les tendances temporelles.',
    'Besoin identifie : ajouter le dataset 2022 et fusionner 2022 + 2023 pour renforcer l historique.',
    'Resultat actuel : 24 maladies modelisees et 20 retenues en production avec R2 >= 0.5.',
    'Message cle : le systeme est defendable en production quand on garde une fenetre de 4 semaines.',
], 0.8, 2.0, 11.7, 3.0)
card(slide, 'Datasets', '2022 + 2023', 0.8, 6.0, 2.7, 0.7, PRIMARY)
card(slide, 'Couverture', '27 maladies | 26 provinces | 517 zones', 3.9, 6.0, 4.8, 0.7, ACCENT)
card(slide, 'Production', '20 modeles eligibles', 9.1, 6.0, 3.0, 0.7, GOOD)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, '1. Pourquoi ajouter le 2eme dataset', 'Difficulte initiale puis justification de la fusion')
bullets(slide, [
    'Le premier dataset seul donnait une profondeur historique trop courte pour plusieurs maladies.',
    'Certaines series etaient trop creuses et rendaient l apprentissage instable.',
    'En reduisant l historique, les performances se degradaient tres vite.',
    'L ajout du dataset 2022 a augmente le nombre de semaines observees et a stabilise l apprentissage.',
    'Conclusion : la fusion 2022 + 2023 etait necessaire pour obtenir des predictions fiables.',
], 0.8, 1.9, 7.0, 4.8)
bullets(slide, [
    'Avant : drc-2023_sem08.xlsx seulement',
    'Apres : drc-2022_sem40.xlsx + drc-2023_sem08.xlsx',
    'Benefice : plus de semaines, plus de contexte, meilleure robustesse',
], 8.2, 2.0, 4.1, 3.2, 16)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, '2. Nettoyage et preparation des donnees', 'Ce qui a ete trouve, supprime, corrige et garde')
bullets(slide, [
    'Donnees trouvees : 2 fichiers bruts, 27 maladies, 26 provinces, 517 zones de sante.',
    'Supprime : colonnes techniques inutiles, lignes sans zone, sans maladie, sans cas, dates invalides, annees hors plage, doublons exacts, valeurs negatives.',
    'Corrige : ages convertis en numerique, NaN ages remplaces par 0, population egale a 0 requalifiee en inconnue.',
    'Traite : outliers plafonnes par IQR x3, maladies trop creuses retirees si plus de 70 pourcent de zeros, petits trous interpol es.',
    'Garde : historique propre par maladie et semaine, puis creation des lags, moyennes mobiles, tendance, volatilite et variables calendaires.',
], 0.75, 1.8, 8.1, 5.3, 16)
card(slide, 'Nettoyage', 'brut -> propre -> agrege -> features', 9.0, 2.1, 3.2, 0.85, PRIMARY)
card(slide, 'Historique', '4 semaines pour predire', 9.0, 3.3, 3.2, 0.85, SECONDARY)
card(slide, 'Filtre', 'maladies faibles ecartees', 9.0, 4.5, 3.2, 0.85, ACCENT)
card(slide, 'Sortie', '20 maladies en production', 9.0, 5.7, 3.2, 0.85, GOOD)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, '3. Metriques et choix des modeles', 'Pourquoi ces metriques et pourquoi ces algorithmes')
bullets(slide, [
    'R2 : mesure principale; plus il est proche de 1, plus le modele explique correctement la realite.',
    'MAE : erreur moyenne absolue; facile a lire en nombre de cas.',
    'RMSE : penalise davantage les grosses erreurs, utile pour les pics epidemiques.',
    'MAPE : erreur relative en pourcentage, pratique quand les volumes changent selon la maladie.',
    'CV R2 : controle de la stabilite dans le temps via validation croisee temporelle.',
    'Plusieurs algorithmes sont compares : Random Forest, Gradient Boosting, Ridge, SVR, KNN, regression lineaire. Le meilleur est retenu par maladie.',
], 0.75, 1.8, 7.8, 5.2, 16)
bullets(slide, [
    'Regle : R2 >= 0.5 pour autoriser la prediction',
    'Exemples >= 0.9 : COVID-19 0.906, Diarrhee aqueuse 0.939, Diarrhee sanglante 0.980, IRA 0.986, Paludisme confirme 0.975',
], 8.8, 2.0, 3.5, 3.6, 15)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, '4. Resultats obtenus par maladie', 'Synthese des performances retenues en production')
bullets(slide, [
    '24 maladies ont ete modelisees; 20 passent le seuil de production (R2 >= 0.5).',
    '16 maladies atteignent un niveau tres fort (R2 >= 0.8).',
    'Maladies ecartees : DECES MATERNELS 0.439, MAPI GRAVES 0.007, RAGE -0.559, PESTE -5.951.',
    'Exemples tres solides : IRA 0.986, Diarrhee sanglante 0.980, Paludisme confirme 0.975, Paludisme suspect 0.967, COVID-19 0.906.',
    'Decision locale : chaque maladie garde son meilleur algorithme, ce qui evite de forcer une seule logique sur tous les profils epidemiologiques.',
], 0.75, 1.8, 7.7, 4.9, 16)
if image_path.exists():
    slide.shapes.add_picture(str(image_path), Inches(8.7), Inches(1.95), width=Inches(3.45))
card(slide, 'Conclusion locale', '20 maladies exploitables en prediction', 8.7, 6.0, 3.45, 0.85, GOOD)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, '5. Metriques globales a presenter', 'Lecture executive des resultats consolides')
bullets(slide, [
    'R2 moyen simple = 0.477 : moyenne brute des 24 maladies. Cette valeur baisse parce que quelques maladies faibles degradent la moyenne generale.',
    'R2 moyen pondere = 0.969 : moyenne tenant compte du volume de cas. Elle montre que les maladies les plus frequentes sont tres bien predites.',
    'MAPE moyen global = 28.8 pourcent : erreur relative moyenne sur les predictions retenues.',
    'F1 Macro = 0.816 : qualite moyenne de classification des niveaux d alerte, maladie par maladie, sans favoriser les maladies les plus volumineuses.',
    'F1 Micro = 0.888 : qualite globale agregee de classification, plus sensible aux maladies avec beaucoup d observations.',
    'Lecture directeur : la performance globale est bonne, mais elle doit toujours etre lue avec les resultats par maladie pour eviter les effets de moyenne.',
], 0.75, 1.8, 8.0, 5.0, 16)
card(slide, 'R2 simple', '0.477', 9.1, 2.0, 3.0, 0.8, PRIMARY)
card(slide, 'R2 pondere', '0.969', 9.1, 3.05, 3.0, 0.8, SECONDARY)
card(slide, 'F1 Macro', '0.816', 9.1, 4.1, 3.0, 0.8, ACCENT)
card(slide, 'F1 Micro', '0.888', 9.1, 5.15, 3.0, 0.8, GOOD)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, '6. Pourquoi garder l approche par maladie', 'Pourquoi ne pas faire un seul modele global')
bullets(slide, [
    'Les maladies n ont pas le meme comportement : volumes, saisonnalites, volatilite et vitesse de propagation changent fortement d une maladie a l autre.',
    'Un modele global unique risque d etre domine par les maladies a tres gros volumes et de mal representer les maladies plus rares ou plus specifiques.',
    'L approche par maladie permet de choisir le meilleur algorithme pour chaque profil epidemiologique au lieu de forcer la meme recette partout.',
    'Elle facilite aussi la gouvernance : on peut exclure uniquement les maladies faibles sans penaliser les maladies tres performantes.',
    'Enfin, elle rend l interpretation plus defendable devant la direction : chaque maladie a son R2, son MAE, son RMSE et son niveau de confiance propre.',
], 0.75, 1.85, 7.95, 4.9, 16)
card(slide, 'Approche retenue', '1 maladie = 1 evaluation = 1 decision', 9.0, 2.05, 3.2, 0.9, PRIMARY)
card(slide, 'Avantage', 'on garde les fortes, on ecarte les faibles', 9.0, 3.25, 3.2, 0.9, ACCENT)
card(slide, 'Impact', 'meilleure robustesse et meilleure defense metier', 9.0, 4.45, 3.2, 0.9, GOOD)
card(slide, 'Message final', 'approche par maladie > modele global unique', 9.0, 5.65, 3.2, 0.9, SECONDARY)

prs.save(str(output_path))
print(output_path)
