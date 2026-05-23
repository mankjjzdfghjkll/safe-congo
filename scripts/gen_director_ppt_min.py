from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


desktop = Path.home() / 'Desktop'
root = Path(__file__).resolve().parent.parent
output_path = desktop / 'Presentation_SAFE_CONGO_Directeur_Demain.pptx'
image_path = root / 'models' / 'evaluation' / 'confusion_matrices_combined.png'

PRIMARY = RGBColor(10, 95, 171)
SECONDARY = RGBColor(73, 172, 239)
TEXT = RGBColor(31, 41, 55)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(217, 119, 6)
GOOD = RGBColor(5, 150, 105)
BG = RGBColor(245, 248, 252)


def background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.5))
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()


def title(slide, text, subtext=''):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.75), Inches(12), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    if subtext:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.35), Inches(12), Inches(0.35))
        p2 = sub.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtext
        r2.font.size = Pt(11)
        r2.font.color.rgb = TEXT


def bullets(slide, items, left, top, width, height, size=17):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.bullet = True
        p.space_after = Pt(7)


def card(slide, heading, value, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = heading
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(20)
    r2.font.bold = True
    r2.font.color.rgb = WHITE


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, 'SAFE CONGO - pipeline data et performances du modele', 'Presentation pour direction | 5 slides maximum')
bullets(slide, [
    'Point de depart : un premier dataset 2023 seul, insuffisant pour stabiliser les tendances temporelles.',
    'Besoin identifie : ajouter le dataset 2022 et fusionner 2022 + 2023 pour renforcer l historique.',
    'Resultat actuel : 24 maladies modelisees et 20 retenues en production avec R2 >= 0.5.',
    'Message cle : le systeme est defendable en production quand on garde une fenetre de 4 semaines.',
], 0.8, 2.0, 11.7, 3.0)
card(slide, 'Datasets', '2022 + 2023', 0.8, 6.0, 2.7, 0.7, PRIMARY)
card(slide, 'Couverture', '27 maladies | 26 provinces | 517 zones', 3.9, 6.0, 4.8, 0.7, ACCENT)
card(slide, 'Production', '20 modeles eligibles', 9.1, 6.0, 3.0, 0.7, GOOD)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, '1. Pourquoi ajouter le 2eme dataset', 'Difficulte initiale puis justification de la fusion')
bullets(slide, [
    'Le premier dataset seul donnait une profondeur historique trop courte pour plusieurs maladies.',
    'Certaines series etaient trop creuses et rendaient l apprentissage instable.',
    'En reduisant l historique, les performances se degradaient tres vite.',
    'L ajout du dataset 2022 a augmente le nombre de semaines observees et a stabilise l apprentissage.',
    'Conclusion : la fusion 2022 + 2023 etait necessaire pour obtenir des predictions fiables.',
], 0.8, 1.9, 7.0, 4.8)
bullets(slide, [
    'Avant : drc-2023_sem08.xlsx seulement',
    'Apres : drc-2022_sem40.xlsx + drc-2023_sem08.xlsx',
    'Benefice : plus de semaines, plus de contexte, meilleure robustesse',
], 8.2, 2.0, 4.1, 3.2, 16)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, '2. Nettoyage et preparation des donnees', 'Ce qui a ete trouve, supprime, corrige et garde')
bullets(slide, [
    'Donnees trouvees : 2 fichiers bruts, 27 maladies, 26 provinces, 517 zones de sante.',
    'Supprime : colonnes techniques inutiles, lignes sans zone, sans maladie, sans cas, dates invalides, annees hors plage, doublons exacts, valeurs negatives.',
    'Corrige : ages convertis en numerique, NaN ages remplaces par 0, population egale a 0 requalifiee en inconnue.',
    'Traite : outliers plafonnes par IQR x3, maladies trop creuses retirees si plus de 70% de zeros, petits trous interpolés.',
    'Garde : historique propre par maladie et semaine, puis creation des lags, moyennes mobiles, tendance, volatilite et variables calendaires.',
], 0.75, 1.8, 8.1, 5.3, 16)
card(slide, 'Nettoyage', 'brut -> propre -> agrege -> features', 9.0, 2.1, 3.2, 0.85, PRIMARY)
card(slide, 'Historique', '4 semaines pour predire', 9.0, 3.3, 3.2, 0.85, SECONDARY)
card(slide, 'Filtre', 'maladies faibles ecartees', 9.0, 4.5, 3.2, 0.85, ACCENT)
card(slide, 'Sortie', '20 maladies en production', 9.0, 5.7, 3.2, 0.85, GOOD)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, '3. Metriques et choix des modeles', 'Pourquoi ces metriques et pourquoi ces algorithmes')
bullets(slide, [
    'R2 : mesure principale; plus il est proche de 1, plus le modele explique correctement la realite.',
    'MAE : erreur moyenne absolue; facile a lire en nombre de cas.',
    'RMSE : penalise davantage les grosses erreurs, utile pour les pics epidemiques.',
    'MAPE : erreur relative en pourcentage, pratique quand les volumes changent selon la maladie.',
    'CV R2 : controle de la stabilite dans le temps via validation croisee temporelle.',
    'Plusieurs algorithmes sont compares : Random Forest, Gradient Boosting, Ridge, SVR, KNN, regression lineaire. Le meilleur est retenu par maladie.',
], 0.75, 1.8, 7.8, 5.2, 16)
bullets(slide, [
    'Regle : R2 >= 0.5 pour autoriser la prediction',
    'Exemples >= 0.9 : COVID-19 0.906, Diarrhee aqueuse 0.939, Diarrhee sanglante 0.980, IRA 0.986, Paludisme confirme 0.975',
], 8.8, 2.0, 3.5, 3.6, 15)

slide = prs.slides.add_slide(prs.slide_layouts[6])
background(slide)
title(slide, '4. Resultats obtenus et message final', 'Synthese pour la decision de direction')
bullets(slide, [
    '24 maladies ont ete modelisees; 20 passent le seuil de production (R2 >= 0.5).',
    '16 maladies atteignent un niveau tres fort (R2 >= 0.8).',
    'Maladies ecartees : DECES MATERNELS 0.439, MAPI GRAVES 0.007, RAGE -0.559, PESTE -5.951.',
    'Performance globale du journal : R2 moyen pondere 96.9%, F1 macro 81.6%, F1 micro 88.8%.',
    'Decision recommandee : conserver la fenetre de 4 semaines, garder le filtre R2 et continuer a alimenter regulierement les donnees terrain.',
], 0.75, 1.8, 7.7, 4.9, 16)
if image_path.exists():
    slide.shapes.add_picture(str(image_path), Inches(8.7), Inches(1.95), width=Inches(3.45))
card(slide, 'Conclusion', 'SAFE CONGO est presentable en mode 4 semaines + filtrage R2', 8.7, 6.0, 3.45, 0.85, GOOD)

prs.save(str(output_path))
print(output_path)