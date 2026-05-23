from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT_DIR = Path(__file__).resolve().parent.parent
DESKTOP_DIR = Path.home() / "Desktop"
OUTPUT_PATH = DESKTOP_DIR / "Presentation_SAFE_CONGO_Directeur_Demain.pptx"
CONFUSION_IMAGE = ROOT_DIR / "models" / "evaluation" / "confusion_matrices_combined.png"

BG = RGBColor(245, 248, 252)
PRIMARY = RGBColor(10, 95, 171)
SECONDARY = RGBColor(73, 172, 239)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(95, 110, 129)
ACCENT = RGBColor(217, 119, 6)
GOOD = RGBColor(5, 150, 105)



def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(13.33),
        Inches(0.55),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.75), Inches(12.2), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = PRIMARY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(1.38), Inches(11.9), Inches(0.45))
        p2 = sub_box.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = MUTED


def add_footer(slide, index_text: str):
    footer = slide.shapes.add_textbox(Inches(10.9), Inches(7.0), Inches(1.8), Inches(0.25))
    p = footer.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = index_text
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)
        p.bullet = True


def add_callout(slide, title: str, value: str, left: float, top: float, width: float, height: float, fill_color: RGBColor):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(23)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(255, 255, 255)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(
        slide,
        "SAFE CONGO - Pipeline data et performances du modele",
        f"Presentation de synthese pour direction | {datetime.now().strftime('%d/%m/%Y')}",
    )
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.95), Inches(12.1), Inches(1.0))
    banner.fill.solid()
    banner.fill.fore_color.rgb = SECONDARY
    banner.line.fill.background()
    tf = banner.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Objectif : montrer comment nous sommes passes d'un dataset insuffisant a un pipeline fiable, traçable et exploitable pour la prediction epidemiologique."
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    add_bullets(
        slide,
        [
            "Point de depart : un premier dataset 2023 seulement, utile mais trop limite pour stabiliser les tendances temporelles.",
            "Decision cle : ajouter le dataset 2022 puis fusionner 2022 + 2023 afin d'augmenter la profondeur historique.",
            "Resultat final actuel : 24 maladies modelisees, 20 maladies retenues en production (R² >= 0.5).",
            "Message directeur : le systeme est plus solide quand il exploite une fenetre de 4 semaines et un filtrage strict des modeles faibles.",
        ],
        0.8,
        3.25,
        11.8,
        2.8,
        18,
    )
    add_callout(slide, "Datasets sources", "2 fichiers : 2022 + 2023", 0.8, 6.15, 2.9, 0.7, PRIMARY)
    add_callout(slide, "Couverture brute", "27 maladies | 26 provinces | 517 zones", 4.05, 6.15, 4.2, 0.7, ACCENT)
    add_callout(slide, "Maladies retenues", "20 modeles eligibles", 8.65, 6.15, 3.5, 0.7, GOOD)
    add_footer(slide, "1 / 5")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "1. Pourquoi le premier dataset ne suffisait pas", "Problemes rencontres puis justification de la fusion 2022 + 2023")
    add_bullets(
        slide,
        [
            "Le premier dataset seul donnait une couverture temporelle trop courte pour certaines maladies et perturbait la lecture des saisonnalites.",
            "Des series etaient trop creuses, certaines maladies avaient trop peu de semaines pour apprendre un signal stable.",
            "Avec un historique trop court, la prediction devient instable et les performances se degradent vite des qu'on reduit la fenetre temporelle.",
            "L'ajout du 2eme dataset a permis d'etendre le nombre de semaines observees, de consolider les signaux et de mieux discriminer les tendances reelles.",
        ],
        0.75,
        1.8,
        7.0,
        4.8,
        18,
    )
    right_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.9), Inches(4.35), Inches(4.8))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_box.line.color.rgb = SECONDARY
    tf = right_box.text_frame
    tf.clear()
    for idx, line in enumerate([
        "Avant fusion",
        "Dataset principal : drc-2023_sem08.xlsx",
        "Limite : profondeur historique insuffisante",
        "Risque : apprentissage fragile sur maladies peu denses",
        "",
        "Apres fusion",
        "Ajout : drc-2022_sem40.xlsx",
        "Gain : plus de semaines, plus de stabilite, meilleur apprentissage",
        "Conclusion : la fusion etait necessaire, pas optionnelle",
    ]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16 if idx in (0,5) else 15)
        p.font.bold = idx in (0,5)
        p.font.color.rgb = PRIMARY if idx in (0,5) else TEXT
        p.space_after = Pt(5)
    add_footer(slide, "2 / 5")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "2. Nettoyage et preparation des donnees", "Ce qui a ete trouve, supprime, corrige et conserve")
    add_bullets(
        slide,
        [
            "Données trouvees : deux fichiers bruts (2022 et 2023), 27 maladies, 26 provinces et 517 zones de sante au niveau geolocalise.",
            "Colonnes inutiles supprimees : variables techniques sans valeur predictive directe.",
            "Lignes ecartees : zone absente, maladie absente, cas manquants, dates invalides, annees hors plage 2020-2030, doublons exacts et valeurs negatives.",
            "Valeurs traitees : ages convertis en numerique, NaN age remplaces par 0, population egale a 0 requalifiee en inconnue.",
            "Series preparees : plafonnement des outliers par IQR x3, suppression des maladies trop creuses (>70% de zeros), interpolation de petits trous isoles.",
            "Donnees gardees pour le modele : historique par maladie/semaine, puis creation des lags, moyennes mobiles, tendance, volatilite et variables calendaires.",
        ],
        0.7,
        1.75,
        8.0,
        5.25,
        17,
    )
    add_callout(slide, "Nettoyage", "brut -> propre -> agrégé -> features", 9.1, 2.0, 3.1, 0.85, PRIMARY)
    add_callout(slide, "Historique retenu", "fenetre de 4 semaines pour predire", 9.1, 3.15, 3.1, 0.85, SECONDARY)
    add_callout(slide, "Filtre qualite", "modeles faibles mis a l'ecart", 9.1, 4.3, 3.1, 0.85, ACCENT)
    add_callout(slide, "Sortie finale", "20 maladies eligibles en production", 9.1, 5.45, 3.1, 0.85, GOOD)
    add_footer(slide, "3 / 5")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "3. Metriques et choix des modeles", "Pourquoi ces indicateurs et comment lire les resultats")
    add_bullets(
        slide,
        [
            "R² : part de la variance expliquee. Plus il est proche de 1, plus le modele suit correctement la realite. C'est la metrique de decision principale.",
            "MAE : erreur moyenne absolue en nombre de cas. Elle montre l'ecart moyen concret entre prevision et reel.",
            "RMSE : penalise davantage les grosses erreurs. Utile quand on veut sanctionner les ratés importants en situation epidemique.",
            "MAPE : erreur relative en pourcentage. Elle facilite la lecture business quand les volumes varient fortement d'une maladie a l'autre.",
            "CV R² : validation croisee temporelle. Elle confirme que le modele tient dans le temps et ne memorise pas seulement le jeu de test final.",
            "Pourquoi plusieurs algorithmes ? Random Forest, Gradient Boosting, Ridge, SVR, KNN et Regression lineaire sont compares; seul le meilleur par maladie est retenu.",
        ],
        0.7,
        1.75,
        7.8,
        5.3,
        16,
    )
    metric_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.85), Inches(3.45), Inches(5.15))
    metric_panel.fill.solid()
    metric_panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    metric_panel.line.color.rgb = PRIMARY
    tf = metric_panel.text_frame
    tf.clear()
    for idx, line in enumerate([
        "Regle de decision",
        "R² >= 0.5 : maladie autorisee en prediction",
        "R² < 0.5 : maladie mise a l'ecart",
        "",
        "Exemples >= 0.9",
        "COVID-19 : 0.906",
        "Diarrhee aqueuse : 0.939",
        "Diarrhee sanglante : 0.980",
        "Paludisme confirme : 0.975",
        "IRA : 0.986",
    ]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16 if idx in (0,4) else 14)
        p.font.bold = idx in (0,4)
        p.font.color.rgb = PRIMARY if idx in (0,4) else TEXT
        p.space_after = Pt(4)
    add_footer(slide, "4 / 5")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "4. Resultats obtenus et message de direction", "Ce qu'il faut retenir pour la decision")
    add_bullets(
        slide,
        [
            "24 maladies ont ete modelisees au total; 20 passent le seuil de production (R² >= 0.5).",
            "16 maladies atteignent meme un niveau tres fort (R² >= 0.8), ce qui confirme la qualite du pipeline en 4 semaines.",
            "Les maladies ecartees sont : DECES MATERNELS (0.439), MAPI GRAVES (0.007), RAGE (-0.559) et PESTE (-5.951).",
            "Performance globale confirmee dans le journal : R² moyen pondere 96.9%, F1 macro 81.6%, F1 micro 88.8%.",
            "Recommendation a valider : conserver la fenetre de 4 semaines, garder le filtrage des modeles faibles et poursuivre l'alimentation reguliere de la base terrain.",
        ],
        0.7,
        1.75,
        7.6,
        4.6,
        16,
    )
    if CONFUSION_IMAGE.exists():
        slide.shapes.add_picture(str(CONFUSION_IMAGE), Inches(8.65), Inches(1.95), width=Inches(3.55))
    add_callout(slide, "Decision finale", "Le systeme est presentable et defendable si on garde la logique 4 semaines + filtrage R².", 8.65, 6.0, 3.55, 0.9, GOOD)
    add_footer(slide, "5 / 5")

    return prs


def main() -> None:
    presentation = build_presentation()
    presentation.save(str(OUTPUT_PATH))
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()